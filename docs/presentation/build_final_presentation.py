from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\redma\Desktop\CapstoneProject")
OUT = ROOT / "docs" / "presentation" / "CSC482_Team2_Final_Presentation.pptx"
SCREENSHOTS = ROOT / "docs" / "screenshots"

NAVY = RGBColor(20, 34, 56)
TEAL = RGBColor(9, 132, 119)
RED = RGBColor(190, 24, 24)
GOLD = RGBColor(188, 117, 25)
LIGHT_BG = RGBColor(244, 247, 251)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(13, 24, 38)
MUTED = RGBColor(82, 99, 123)
BORDER = RGBColor(210, 219, 232)


def add_fill(slide, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def textbox(slide, text, x, y, w, h, size=24, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = "Aptos"
    p.font.color.rgb = color
    return box


def title(slide, text, subtitle=None):
    textbox(slide, text, 0.7, 0.45, 11.9, 0.55, size=30, bold=True, color=NAVY)
    if subtitle:
        textbox(slide, subtitle, 0.72, 1.04, 11.9, 0.35, size=12, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.42), Inches(11.9), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def bullet_list(slide, items, x, y, w, h, size=17, color=DARK, gap=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def card(slide, x, y, w, h, heading, body=None, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    textbox(slide, heading, x + 0.2, y + 0.16, w - 0.35, 0.35, size=16, bold=True, color=NAVY)
    if body:
        if isinstance(body, list):
            bullet_list(slide, body, x + 0.2, y + 0.62, w - 0.35, h - 0.75, size=12, color=DARK)
        else:
            textbox(slide, body, x + 0.2, y + 0.62, w - 0.35, h - 0.75, size=12, color=DARK)
    return shape


def image_fit(slide, image_path, x, y, w, h):
    path = Path(image_path)
    if not path.exists():
        card(slide, x, y, w, h, "Screenshot placeholder", str(path.name), accent=GOLD)
        return
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pic_w = iw * scale
    pic_h = ih * scale
    left = x + (w - pic_w) / 2
    top = y + (h - pic_h) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(pic_w), height=Inches(pic_h))
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    border.fill.background()
    border.line.color.rgb = BORDER


def header_bar(slide, label="CSC 482 Capstone Project 2 - Team 2"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    textbox(slide, label, 0.7, 0.08, 7, 0.2, size=9, color=WHITE)


def add_slide_base(prs, heading, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_fill(slide, LIGHT_BG)
    header_bar(slide)
    title(slide, heading, subtitle)
    return slide


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_fill(slide, NAVY)
    textbox(slide, "Intelligent Security Log Summarization\nand Threat Timeline Generation", 0.8, 1.45, 11.7, 1.45, size=36, bold=True, color=WHITE)
    textbox(slide, "CSC 482 Capstone Project 2 | Project Team 2", 0.83, 3.05, 10.5, 0.38, size=18, color=WHITE)
    textbox(slide, "GitHub: https://github.com/Dredman72/security-log-threat-timeline", 0.83, 3.55, 10.8, 0.35, size=14, color=RGBColor(196, 230, 255))
    card(slide, 0.85, 4.55, 3.55, 1.2, "Derrick Redman", "Team Lead + AI/LLM Specialist")
    card(slide, 4.9, 4.55, 3.55, 1.2, "Zion Moore", "Backend & Log Processing Lead")
    card(slide, 8.95, 4.55, 3.55, 1.2, "Oriah Molton-Bowman", "Frontend & Visualization Lead")

    # 2
    slide = add_slide_base(prs, "Presentation Plan", "15 minute presentation + 8 minute demo + 5 minute Q&A")
    card(slide, 0.8, 1.8, 3.7, 3.7, "15 Minute Presentation", [
        "Problem and project goal",
        "Architecture and workflow",
        "Team contributions",
        "Testing, lessons learned, future work",
    ])
    card(slide, 4.85, 1.8, 3.7, 3.7, "8 Minute Demo", [
        "Open Flask app",
        "Paste or upload security logs",
        "Review generated report",
        "Download JSON and HTML report",
    ], accent=GOLD)
    card(slide, 8.9, 1.8, 3.7, 3.7, "5 Minute Q&A", [
        "Explain design choices",
        "Discuss limitations",
        "Answer deployment and testing questions",
    ], accent=RED)

    # 3
    slide = add_slide_base(prs, "Problem and Motivation")
    bullet_list(slide, [
        "Security analysts often review large volumes of Windows, firewall, authentication, Sysmon, and EDR logs.",
        "Important events can be buried in repeated failures, noisy network activity, and raw technical details.",
        "Manual review makes it difficult to quickly explain what happened, when it happened, and why it matters.",
        "Our prototype helps convert raw log evidence into a readable security report and chronological threat timeline.",
    ], 0.95, 1.9, 6.2, 4.5, size=19, gap=8)
    card(slide, 7.6, 2.1, 4.8, 2.9, "Target User", [
        "Blue team analyst",
        "SOC student or beginner analyst",
        "Incident response reviewer",
        "Instructor reviewing log analysis output",
    ])

    # 4
    slide = add_slide_base(prs, "Project Goal and Scope")
    card(slide, 0.85, 1.8, 5.65, 3.85, "Primary Goal", "Build a local Flask application that accepts uploaded or pasted security logs, uses OpenAI to summarize likely threats, and generates a chronological timeline with evidence.")
    card(slide, 6.85, 1.8, 5.65, 3.85, "Scope", [
        "No model training required",
        "Prototype runs locally",
        "Deployable on Windows Server 2025 and Ubuntu 24.04 LTS",
        "Focus on structured output, timeline generation, and report clarity",
    ], accent=GOLD)

    # 5
    slide = add_slide_base(prs, "System Workflow")
    steps = [
        ("1. Input", "Paste logs or upload .txt, .log, .csv, .json, or .xml files"),
        ("2. Parse", "Normalize events into fields such as timestamp, source, user, IP, event type, severity, and evidence"),
        ("3. Analyze", "OpenAI returns structured JSON with summary, risk, attack type, timeline, and actions"),
        ("4. Render", "Flask displays the report, parsed event preview, timeline, and download options"),
    ]
    x = 0.7
    for i, (h, b) in enumerate(steps):
        card(slide, x + i * 3.15, 2.1, 2.75, 2.7, h, b, accent=[TEAL, GOLD, RED, NAVY][i])
    textbox(slide, "End-to-end result: raw evidence becomes a report-ready incident narrative.", 1.25, 5.55, 10.8, 0.5, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 6
    slide = add_slide_base(prs, "Current Prototype Features")
    bullet_list(slide, [
        "Flask web interface for log input and upload.",
        "OpenAI-powered structured security summary.",
        "Risk level, risk rationale, attack type, assets, indicators, findings, and recommended actions.",
        "Threat timeline with timestamp, source, severity, details, and evidence.",
        "Parsed event preview and downloadable JSON/HTML report output.",
    ], 0.75, 1.7, 4.4, 4.8, size=15, gap=5)
    image_fit(slide, SCREENSHOTS / "Flask App Homepage.png", 5.35, 1.75, 7.25, 4.7)

    # 7
    slide = add_slide_base(prs, "Derrick: LLM Prompt and Structured Output")
    card(slide, 0.8, 1.75, 5.8, 4.6, "Structured JSON Report Fields", [
        "executive_summary and generated_at",
        "risk_level and risk_rationale",
        "attack_type",
        "affected_assets and indicators_of_compromise",
        "key_findings",
        "timeline with evidence",
        "recommended_actions",
    ])
    card(slide, 6.95, 1.75, 5.55, 4.6, "Quality Rules", [
        "Return only valid JSON",
        "Do not invent unsupported events",
        "Cite evidence from logs",
        "Use concise language",
        "Handle imperfect output with fallback behavior",
    ], accent=RED)

    # 8
    slide = add_slide_base(prs, "Threat Timeline Generation")
    bullet_list(slide, [
        "Timeline events must include timestamp, source, event, severity, details, and evidence.",
        "Repeated similar events should be grouped when useful.",
        "Event titles should stay short and readable.",
        "High-impact events such as successful logins, privilege escalation, sensitive file access, and firewall blocks are prioritized.",
    ], 0.75, 1.75, 4.6, 4.8, size=15, gap=5)
    image_fit(slide, SCREENSHOTS / "OpenAI Report Output2.png", 5.55, 1.72, 6.95, 4.75)

    # 9
    slide = add_slide_base(prs, "Zion: Backend Parser and Data Handoff")
    card(slide, 0.75, 1.7, 4.7, 4.8, "Backend Responsibilities", [
        "Canonical event schema",
        "Sample log parser testing",
        "Normalized JSON output",
        "Parser-to-LLM input format",
        "Upload-to-report pipeline notes",
    ])
    image_fit(slide, SCREENSHOTS / "Updated Normalized Events JSON.png", 5.7, 1.7, 6.8, 4.8)

    # 10
    slide = add_slide_base(prs, "Oriah: Frontend and Visualization")
    card(slide, 0.75, 1.7, 4.7, 4.8, "Frontend Responsibilities", [
        "Report wireframe",
        "HTML report layout",
        "Timeline visualization plan",
        "Report data field documentation",
        "Readable report sections for demo review",
    ], accent=GOLD)
    image_fit(slide, SCREENSHOTS / "Wirefram Powerpoint.png", 5.7, 1.7, 6.8, 4.8)

    # 11
    slide = add_slide_base(prs, "Testing and Validation")
    card(slide, 0.75, 1.65, 3.75, 4.95, "Sample Scenarios", [
        "SSH brute-force/root access",
        "Firewall block only",
        "Suspicious PowerShell",
        "EDR malware alert",
        "Mixed authentication, firewall, and EDR activity",
    ])
    card(slide, 4.8, 1.65, 3.75, 4.95, "Validation Checks", [
        "Required parser fields",
        "Prompt output structure",
        "Timeline evidence",
        "Fallback behavior",
        "Downloadable JSON and HTML output",
    ], accent=RED)
    card(slide, 8.85, 1.65, 3.75, 4.95, "Result", [
        "Prototype runs locally",
        "Reports are readable",
        "Timeline supports incident review",
        "Outputs support documentation and screenshots",
    ], accent=TEAL)

    # 12
    slide = add_slide_base(prs, "Milestones Achieved")
    rows = [
        ("Week 1", "Lab setup, requirements, OpenAI API plan"),
        ("Week 2", "Initial prototypes, JSON output structure, parser schema, wireframe"),
        ("Week 3", "Parser validation checklist and early parser testing"),
        ("Week 4", "OpenAI integration, structured summaries, fallback handling"),
        ("Week 5", "Threat timeline rules, event grouping, evidence requirements"),
        ("Week 6", "End-to-end report generation, reset, JSON/HTML downloads"),
        ("Week 7", "Integration testing, prompt consistency, multi-scenario logs"),
        ("Week 8", "Final report, final presentation, final demo preparation"),
    ]
    for idx, (wk, desc) in enumerate(rows):
        y = 1.65 + idx * 0.55
        textbox(slide, wk, 1.0, y, 1.1, 0.28, size=13, bold=True, color=TEAL)
        textbox(slide, desc, 2.2, y, 9.7, 0.28, size=13, color=DARK)

    # 13
    slide = add_slide_base(prs, "Team Contributions")
    card(slide, 0.75, 1.75, 3.85, 4.75, "Derrick Redman", [
        "Team leadership",
        "LLM prompt and JSON schema",
        "Risk rationale and timeline rules",
        "Report download/reset workflow",
        "Testing documentation",
    ], accent=TEAL)
    card(slide, 4.75, 1.75, 3.85, 4.75, "Zion Moore", [
        "Backend parser schema",
        "Sample logs",
        "Normalized event output",
        "Parser-to-report workflow",
        "Integration checks",
    ], accent=GOLD)
    card(slide, 8.75, 1.75, 3.85, 4.75, "Oriah Molton-Bowman", [
        "Report wireframe",
        "Timeline layout planning",
        "HTML report generator",
        "Frontend report fields",
        "Visualization support",
    ], accent=RED)

    # 14
    slide = add_slide_base(prs, "Demo Plan")
    bullet_list(slide, [
        "1. Open the local Flask application.",
        "2. Paste or upload a sample security log file.",
        "3. Click Analyze Logs.",
        "4. Review the executive summary, risk level, attack type, and risk rationale.",
        "5. Walk through assets, indicators, key findings, timeline, and evidence.",
        "6. Show parsed event preview.",
        "7. Download JSON and HTML report outputs.",
        "8. Use Reset to prepare for another scan.",
    ], 0.9, 1.75, 5.4, 4.9, size=16, gap=3)
    image_fit(slide, SCREENSHOTS / "Week 6 HTML Report Generator.png", 6.55, 1.75, 5.85, 4.9)

    # 15
    slide = add_slide_base(prs, "Lessons Learned and Future Work")
    card(slide, 0.75, 1.75, 5.7, 4.75, "Lessons Learned", [
        "Structured JSON makes the frontend and backend easier to connect.",
        "Timeline evidence is critical for trust.",
        "Parser fields must be consistent across sample log types.",
        "Downloadable reports make testing and submission easier.",
    ])
    card(slide, 6.85, 1.75, 5.7, 4.75, "Future Work", [
        "Add broader SIEM and cloud log support.",
        "Improve source IP and destination IP extraction.",
        "Add user authentication and deployment hardening.",
        "Add threat intelligence enrichment and stronger validation.",
    ], accent=GOLD)
    textbox(slide, "Q&A", 5.75, 6.55, 1.8, 0.45, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(make_deck())
